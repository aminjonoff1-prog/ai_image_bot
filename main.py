import os
import asyncio
import logging
import time
import re

from aiohttp import web
from aiogram import Bot, Dispatcher, F
from aiogram.types import Message, FSInputFile, BotCommand
from aiogram.filters import Command
from aiogram.client.default import DefaultBotProperties
from aiogram.enums import ParseMode

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import google.generativeai as genai

# --- CONFIG IMPORT ---
try:
    from config import BOT_TOKEN, FREE_LIMIT, GEMINI_API_KEY, ADMIN_ID
except ImportError:
    BOT_TOKEN = os.environ.get("BOT_TOKEN", "")
    FREE_LIMIT = int(os.environ.get("FREE_LIMIT", "5"))
    GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")
    ADMIN_ID = int(os.environ.get("ADMIN_ID", "0"))

# --- DB FUNKSIYALARI ---
from db import (
    check_limit,
    init_db,
    add_user,
    get_stats,
    get_all_users,
    add_premium_limit,
    get_limit_info,
    reset_user_usage,
    get_user_info,
)

# --- MODULLAR ---
from keyboards import main_menu
from image_gen import generate_image
from video_gen import generate_video
from audio_gen import generate_audio
from coursework_gen import generate_coursework

# --- LOGGING ---
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

# --- GEMINI SOZLAMALARI ---
MODEL_NAMES = [
    "gemini-2.0-flash",
    "gemini-1.5-flash",
    "gemini-1.5-pro",
]

primary_gemini_model = None
if GEMINI_API_KEY:
    try:
        genai.configure(api_key=GEMINI_API_KEY.strip())
        primary_gemini_model = genai.GenerativeModel(MODEL_NAMES[0])
        logger.info("✅ Gemini muvaffaqiyatli yuklandi")
    except Exception as e:
        logger.warning(f"⚠️ Gemini init xatosi: {e}")
        primary_gemini_model = None

# --- BOT VA DISPATCHER ---
bot = Bot(
    token=BOT_TOKEN.strip(),
    default=DefaultBotProperties(parse_mode=ParseMode.HTML),
)
dp = Dispatcher()

# --- FOYDALANUVCHI HOLATLARI ---
user_category = {}

# ============================================================
# YORDAMCHI FUNKSIYALAR (Helper Functions)
# ============================================================

def is_admin(user_id: int) -> bool:
    return user_id == ADMIN_ID


async def safe_delete(message):
    try:
        await message.delete()
    except Exception:
        pass


async def safe_remove_file(filename: str):
    try:
        if filename and os.path.exists(filename):
            os.remove(filename)
    except Exception as e:
        logger.error(f"Faylni o'chirishda xato: {e}")


# --- DIZAYN FUNKSIYALARI ---

def set_slide_background(slide, color1, color2=None):
    """Slaydga qattiq rang yoki gradient beradi."""
    fill = slide.background.fill
    if color2:
        fill.gradient()
        fill.gradient_angle = 90
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2
    else:
        fill.solid()
        fill.fore_color.rgb = color1


def add_title(slide, title_text, subtitle_text="", dark=True):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title_text
    p1.font.name = "Georgia"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(11, 29, 58) if dark else RGBColor(255, 255, 255)
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = "Calibri"
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0, 102, 204) if dark else RGBColor(212, 175, 55)
        p2.space_before = Pt(3)


def add_footer(slide, topic, num):
    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(8.8), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Loyiha: {topic[:35]}... | Slayd {num}/25 | Maxfiy"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(130, 130, 130)


def add_bullets_block(slide, bullets, left=0.7, top=1.55, width=5.0, height=4.6):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:5]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(40, 50, 65)
        p.space_after = Pt(10)
        p.line_spacing = 1.12


def add_callout_card(slide, title, text, left=6.1, top=1.65, width=3.1, height=4.5):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(0, 102, 204)
    card.line.width = Pt(1.2)
    box = slide.shapes.add_textbox(Inches(left + 0.12), Inches(top + 0.15), Inches(width - 0.24), Inches(height - 0.3))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Calibri"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0, 102, 204)
    p1.space_after = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.name = "Georgia"
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(11, 29, 58)
    p2.line_spacing = 1.18


# ============================================================
# AI VA LOGIKA (Timeout va Xatoliklarga qarshi)
# ============================================================

async def gemini_generate_text(prompt: str) -> str:
    """
    Xavfsiz generatsiya: 60 sekund kutadi, timeout bo'lsa qayta urinadi.
    """
    if not GEMINI_API_KEY:
        raise RuntimeError("API KEY yo'q")

    # 1. Primary model
    if primary_gemini_model:
        for attempt in range(2): # 2 marta urinadi
            try:
                response = await asyncio.wait_for(
                    asyncio.to_thread(primary_gemini_model.generate_content, prompt),
                    timeout=60.0  # 60 sekund limit
                )
                text = (response.text or "").strip()
                if text: return text
            except asyncio.TimeoutError:
                logger.warning(f"⏱️ Timeout (urinish {attempt+1})")
                if attempt == 1: break # Ikkinchi urinish ham bo'ldi
                await asyncio.sleep(1)
            except Exception as e:
                logger.warning(f"Model xatosi: {e}")
                break

    # 2. Fallback modellar
    for model_name in MODEL_NAMES:
        try:
            model = genai.GenerativeModel(model_name)
            response = await asyncio.wait_for(
                asyncio.to_thread(model.generate_content, prompt),
                timeout=60.0
            )
            text = (response.text or "").strip()
            if text: return text
        except Exception:
            continue

    raise RuntimeError("AI javob bera olmadi")


async def get_ai_slide_content(topic, slide_num):
    """Har bir slayd uchun kontent olish (xavfsiz)."""
    # Default ma'lumotlar (AI ishlamay qolsa shular qo'yiladi)
    default_data = {
        "title": f"Slayd {slide_num}",
        "subtitle": f"{topic} bo'yicha tahlil",
        "bullets": ["Mavzu tahlili", "Strategik yondashuv", "Tavsiyalar", "Xulosa"],
        "callout": "Tizim tomonidan avtomatik tahlil."
    }

    if not GEMINI_API_KEY:
        return default_data

    prompt = f"""
MAVZU: {topic}
SLAYD: {slide_num}
QISQA FORMATDA JAVOB BERING:
TITLE: [Sarlavha]
SUBTITLE: [Pastki sarlavha]
BULLETS:
- [Nuqta 1]
- [Nuqta 2]
- [Nuqta 3]
CALLOUT: [Qisqa xulosa]
"""

    try:
        text = await gemini_generate_text(prompt)
        data = {"title": "", "subtitle": "", "bullets": [], "callout": ""}
        current_section = None
        
        for line in text.splitlines():
            line = line.strip().upper()
            if line.startswith("TITLE:"): data["title"] = line.split(":",1)[1].strip()
            elif line.startswith("SUBTITLE:"): data["subtitle"] = line.split(":",1)[1].strip()
            elif line.startswith("BULLETS:"): current_section = "bullets"
            elif line.startswith("CALLOUT:"): 
                data["callout"] = line.split(":",1)[1].strip()
                current_section = None
            elif current_section == "bullets":
                clean = re.sub(r"^[-•*\d\.\)\s]+", "", line).strip()
                if clean: data["bullets"].append(clean)

        # Validatsiya
        if not data["title"]: data["title"] = default_data["title"]
        if not data["bullets"]: data["bullets"] = default_data["bullets"]
        return data

    except Exception as e:
        logger.error(f"Slayd {slide_num} generatsiyasi xato: {e}")
        return default_data


async def add_image_to_right_panel(slide, topic, slide_num):
    """Rasm qo'shish yoki placeholder qo'yish."""
    prompt = f"{topic}, professional minimalist illustration, vector art, white background"
    img_added = False
    try:
        img_file, error = await generate_image(prompt, "🖼 Realistik")
        if img_file and os.path.exists(img_file):
            slide.shapes.add_picture(img_file, Inches(5.9), Inches(1.6), width=Inches(3.8), height=Inches(4.0))
            img_added = True
            await safe_remove_file(img_file)
    except Exception as e:
        logger.error(f"Rasm xatosi: {e}")

    if not img_added:
        # Placeholder
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(1.6), Inches(3.6), Inches(4.0))
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(240, 244, 248)
        ph.line.color.rgb = RGBColor(0, 102, 204); ph.line.width = Pt(2)
        tb = slide.shapes.add_textbox(Inches(6.2), Inches(3.0), Inches(3.2), Inches(1.2))
        tf = tb.text_frame; p = tf.paragraphs[0]
        p.text = "🖼️ VISUAL CONCEPT"
        p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    return img_added


# ============================================================
# PREZENTATSIYA YARATISH (XATOLIKLARGA CHIDAMLI)
# ============================================================

async def create_presentation(topic: str, user_id: int, status_msg=None) -> str:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    
    dark_navy = RGBColor(11, 29, 58)
    light_blue = RGBColor(235, 242, 250)
    white = RGBColor(255, 255, 255)
    accent_blue = RGBColor(0, 102, 204)
    gold = RGBColor(212, 175, 55)

    # 1. COVER
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, dark_navy, RGBColor(40, 60, 90))
    # ... (Cover dizayni kodlari o'zgarishsiz) ...
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.65), Inches(10), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent_blue; bar.line.fill.background
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(8.4), Inches(1.7))
    tf = title_box.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.text = topic.upper(); p1.font.name = "Georgia"; p1.font.size = Pt(30); p1.font.bold = True; p1.font.color.rgb = white
    p2 = tf.add_paragraph(); p2.text = "STRATEGIK TAHLILIY PREZENTATSIYA"; p2.font.name = "Calibri"; p2.font.size = Pt(15); p2.font.bold = True; p2.font.color.rgb = gold; p2.space_before = Pt(16)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.1), Inches(4.9), Inches(5.8), Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb = white; badge.line.fill.background
    badge_text = slide.shapes.add_textbox(Inches(2.2), Inches(5.03), Inches(5.6), Inches(0.3))
    tfb = badge_text.text_frame; pb = tfb.paragraphs[0]
    pb.text = "25 ta premium slayd | AI Generated"; pb.font.name = "Calibri"; pb.font.size = Pt(12); pb.font.bold = True; pb.font.color.rgb = dark_navy; pb.alignment = PP_ALIGN.CENTER

    # 2. AGENDA
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, white, light_blue)
    add_title(slide, "Mundarija", "Asosiy bo'limlar")
    add_bullets_block(slide, ["Dolzarblik", "Tahlil", "Strategiya", "Xulosa"], left=0.9, top=1.7, width=5.0, height=4.5)
    add_callout_card(slide, "MAQSAD", "Mavzuni chuqur o'rganish.", left=6.15, top=1.7, width=3.05, height=4.5)
    add_footer(slide, topic, 2)

    # 3-25 LOOP (Har bir slayd alohida himoyalangan)
    for num in range(3, 26):
        slide = prs.slides.add_slide(blank)
        
        # Fon
        if num % 3 == 0: set_slide_background(slide, white, RGBColor(245, 245, 255))
        elif num % 3 == 1: set_slide_background(slide, RGBColor(248, 250, 252), white)
        else: set_slide_background(slide, white)

        # AI CONTENT (TRY-EXCEPT BLOCK)
        s_data = None
        try:
            s_data = await get_ai_slide_content(topic, num)
        except Exception as e:
            logger.error(f"CRITICAL: Slayd {num} yaratib bo'lmadi: {e}")
            s_data = {"title": f"Slayd {num} (Xatolik)", "subtitle": "", "bullets": ["Ma'lumot olinmadi"], "callout": ""}

        add_title(slide, s_data.get("title", f"Slayd {num}"), s_data.get("subtitle", ""))
        bullets = s_data.get("bullets", [])
        callout = s_data.get("callout", "")

        # LAYOUT 1
        if num % 2 == 0:
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.0), Inches(5.0))
            tf = tb.text_frame; tf.word_wrap = True
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {b}"; p.font.size = Pt(16); p.font.color.rgb = RGBColor(50,50,50); p.space_after = Pt(12)
            
            if not await add_image_to_right_panel(slide, topic, num):
                add_callout_card(slide, "XULOSA", callout, left=5.8, top=1.6, width=3.6, height=4.0)
        
        # LAYOUT 2
        else:
            cb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.0))
            cb.fill.solid(); cb.fill.fore_color.rgb = white
            cb.line.color.rgb = accent_blue; cb.line.width = Pt(1)
            itf = cb.text_frame; itf.word_wrap = True
            itf.margin_top = Inches(0.2); itf.margin_bottom = Inches(0.2); itf.margin_left = Inches(0.2); itf.margin_right = Inches(0.2)
            pt = itf.paragraphs[0]; pt.text = "TAHLIL VA STATISTIKA"
            pt.font.bold = True; pt.font.size = Pt(18); pt.font.color.rgb = accent_blue; pt.alignment = PP_ALIGN.CENTER
            for i, b in enumerate(bullets):
                p = itf.add_paragraph(); p.text = b; p.font.size = Pt(16); p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(15)
            pf = itf.add_paragraph(); pf.text = f"— {callout} —"
            pf.font.italic = True; pf.font.size = Pt(14); pf.font.color.rgb = RGBColor(100,100,100); pf.alignment = PP_ALIGN.CENTER

        add_footer(slide, topic, num)

    filename = f"prezentatsiya_{user_id}_{int(time.time())}.pptx"
    prs.save(filename)
    return filename


# ============================================================
# HANDLERLAR (Qisqartirilgan, lekin to'liq)
# ============================================================

@dp.message(Command("start"))
async def start_command(m: Message):
    add_user(m.from_user.id, m.from_user.username, m.from_user.full_name)
    await m.answer(f"👋 Salom, <b>{m.from_user.full_name}</b>!\n\n🤖 Men AI professional dizaynerman", reply_markup=main_menu())

@dp.message(F.text.in_(["📊 Prezentatsiya", "🎨 Logo", "📚 Kurs Ishi"]))
async def category_handler(m: Message):
    user_category[m.from_user.id] = m.text
    await m.answer(f"✍️ <b>{m.text}</b> uchun mavzu yozing:")

@dp.message()
async def generate_handler(m: Message):
    user_id = m.from_user.id
    if user_id not in user_category:
        await m.answer("❗ Avval kategoriyani tanlang.")
        return
    
    user_text = m.text.strip()
    if not user_text: return

    if not is_admin(user_id) and not check_limit(user_id, FREE_LIMIT):
        await m.answer("❌ Limit tugadi. Admin bilan bog'laning.")
        return

    category = user_category.pop(user_id)
    
    # === PREZENTATSIYA ===
    if category == "📊 Prezentatsiya":
        msg = await m.answer("⏳ 25 ta slayd tayyorlanmoqda... (Iltimos, kutib turing)")
        try:
            file_path = await create_presentation(user_text, user_id)
            if os.path.exists(file_path):
                await m.answer_document(FSInputFile(file_path), caption=f"📊 Tayyor! Mavzu: {user_text}")
                await safe_delete(msg)
            else:
                await m.answer("❌ Fayl yaratilmadi.")
        except Exception as e:
            logger.error(f"Prezentatsiya xatosi: {e}")
            await m.answer(f"❌ Xatolik: {e}")
        finally:
            await safe_remove_file(file_path if 'file_path' in locals() else None)
            await safe_delete(msg)

    # BOSHQA KATEGORIYALAR (Mantiq o'zgarmagan)
    elif category == "🎬 Video Generatsiya":
        msg = await m.answer("⏳ Video yuklanmoqda...")
        try:
            f, err = await generate_video(user_text)
            if f and os.path.exists(f): await m.answer_video(FSInputFile(f), caption="🎬 Video tayyor")
            else: await m.answer(f"❌ Xatolik: {err}")
        except Exception as e: await m.answer(f"❌ Xato: {e}")
        finally: 
            await safe_remove_file(f if 'f' in locals() else None)
            await safe_delete(msg)

    elif category == "🎵 Audio/Musiqa":
        msg = await m.answer("⏳ Audio yaratilmoqda...")
        try:
            f, err = await generate_audio(user_text)
            if f and os.path.exists(f): await m.answer_audio(FSInputFile(f), caption="🎵 Audio tayyor")
            else: await m.answer(f"❌ Xatolik: {err}")
        except Exception as e: await m.answer(f"❌ Xato: {e}")
        finally: 
            await safe_remove_file(f if 'f' in locals() else None)
            await safe_delete(msg)
            
    else: # Rasm va boshqalar
        msg = await m.answer("⏳ Rasm chizilmoqda...")
        try:
            f, err = await generate_image(user_text, category)
            if f and os.path.exists(f): await m.answer_photo(FSInputFile(f), caption=f"🎨 {category}")
            else: await m.answer(f"❌ Xatolik: {err}")
        except Exception as e: await m.answer(f"❌ Xato: {e}")
        finally: 
            await safe_remove_file(f if 'f' in locals() else None)
            await safe_delete(msg)

# WEB SERVER VA MAIN
async def run_web_server():
    app = web.Application()
    async def hp(r): return web.Response(text="OK")
    app.router.add_get("/", hp)
    runner = web.AppRunner(app)
    await runner.setup()
    site = web.TCPSite(runner, "0.0.0.0", int(os.environ.get("PORT", 8080)))
    await site.start()

async def main():
    init_db()
    if not BOT_TOKEN: return
    await bot.delete_webhook(drop_pending_updates=True)
    await bot.set_my_commands([BotCommand(command="start", description="Start")])
    logger.info("Bot ishga tushdi.")
    await asyncio.gather(run_web_server(), dp.start_polling(bot))

if __name__ == "__main__":
    try: asyncio.run(main())
    except KeyboardInterrupt: pass
