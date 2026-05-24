import os
import asyncio
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import google.generativeai as genai

from config import GEMINI_API_KEY

# GEMINI
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel("gemini-1.5-flash")

# PROFESSIONAL RANGLAR
COLORS = {
    "primary": RGBColor(15, 23, 42),
    "secondary": RGBColor(30, 41, 59),
    "accent": RGBColor(37, 99, 235),
    "light": RGBColor(248, 250, 252),
    "gray": RGBColor(100, 116, 139)
}


# =========================
# AI SLIDE CONTENT
# =========================

async def generate_slide_content(topic, slide_number):

    prompt = f"""
You are a world-class management consultant and presentation designer.

Create PROFESSIONAL slide content in Uzbek language.

TOPIC:
{topic}

SLIDE NUMBER:
{slide_number}

STRICT RULES:
- Use premium consulting style
- Sound like McKinsey / BCG
- Academic Uzbek language
- Powerful short bullet points
- Include analytical insights
- Include trends and strategic conclusions
- No introductions
- No explanations
- Only slide-ready content

FORMAT:

TITLE:
...

POINTS:
• ...
• ...
• ...
• ...

INSIGHT:
...

Everything in Uzbek.
"""

    try:
        response = await asyncio.to_thread(
            model.generate_content,
            prompt
        )

        return response.text.strip()

    except Exception as e:
        return f"""
TITLE:
{topic}

POINTS:
• Analitik ma'lumot
• Strategik tahlil
• Muhim ko'rsatkichlar
• Zamonaviy tendensiyalar

INSIGHT:
Strategik rivojlanish muhim hisoblanadi.
"""


# =========================
# SLIDE DESIGN
# =========================

def set_slide_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["light"]


def add_title(slide, text):

    title_box = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(0.4),
        Inches(11),
        Inches(0.8)
    )

    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text

    p.font.size = Pt(30)
    p.font.bold = True
    p.font.name = "Aptos"
    p.font.color.rgb = COLORS["primary"]

    p.alignment = PP_ALIGN.LEFT


def add_footer(slide, topic, slide_num):

    footer = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(7.0),
        Inches(12),
        Inches(0.3)
    )

    tf = footer.text_frame
    p = tf.paragraphs[0]

    p.text = f"{topic[:40]} | {slide_num}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["gray"]


def add_content(slide, content):

    content_box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.4),
        Inches(11),
        Inches(5)
    )

    tf = content_box.text_frame

    lines = content.split("\n")

    for line in lines:

        if not line.strip():
            continue

        p = tf.add_paragraph()

        p.text = line.replace("•", "▪")

        if "TITLE:" in line:
            p.font.bold = True
            p.font.size = Pt(24)
            p.font.color.rgb = COLORS["accent"]

        elif "INSIGHT:" in line:
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS["secondary"]

        else:
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS["secondary"]

        p.font.name = "Aptos"

        p.space_after = Pt(8)


# =========================
# MAIN PRESENTATION
# =========================

async def create_professional_presentation(topic, user_id):

    prs = Presentation()

    # TITLE SLIDE
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_slide_background(slide)

    title = slide.shapes.add_textbox(
        Inches(1),
        Inches(2),
        Inches(10),
        Inches(1)
    )

    tf = title.text_frame
    p = tf.paragraphs[0]

    p.text = topic.upper()

    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    subtitle = slide.shapes.add_textbox(
        Inches(1),
        Inches(3),
        Inches(8),
        Inches(1)
    )

    sub_tf = subtitle.text_frame
    sub_p = sub_tf.paragraphs[0]

    sub_p.text = "Professional AI Presentation"
    sub_p.font.size = Pt(20)
    sub_p.font.color.rgb = COLORS["gray"]

    # CONTENT SLIDES
    for i in range(2, 26):

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        set_slide_background(slide)

        content = await generate_slide_content(topic, i)

        add_title(slide, f"{i}-BO'LIM")

        add_content(slide, content)

        add_footer(slide, topic, i)

    filename = f"professional_presentation_{user_id}_{int(time.time())}.pptx"

    prs.save(filename)

    return filename
